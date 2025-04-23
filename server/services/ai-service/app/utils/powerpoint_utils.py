import requests
from langchain.schema import Document
from dotenv import load_dotenv
import base64
import logging
from typing import Dict, List, Optional, Union, Tuple
from pydantic import BaseModel
import tempfile
import os
from azure.storage.blob import BlobServiceClient
from urllib.parse import urlparse
from pptx import Presentation
import asyncio

# Configure logging
logger = logging.getLogger(__name__)



# Load environment variables
load_dotenv()

# Set up logging
logger = logging.getLogger(__name__)


def load_pptx_from_storage_url(
        storage_url: str,
        connection_string: str = None
) -> Optional[List[Document]]:
    """
    Loads a PowerPoint file directly from an Azure Storage URL using connection string.

    Args:
        storage_url: Complete Azure Storage URL to the PowerPoint file
        connection_string: Azure Storage connection string (uses .env if not provided)

    Returns:
        List[Document]: List of Document objects or None if loading fails
    """
    try:
        # Use environment variables if not provided
        if connection_string is None:
            connection_string = os.getenv("AZURE_STORAGE_CONNECTION_STRING")
            if not connection_string:
                logger.error("Azure connection string not provided and not found in environment variables")
                return None

        # Parse the storage URL
        parsed_url = urlparse(storage_url)

        # Check if it's a valid Azure blob storage URL
        if not parsed_url.netloc.endswith('.blob.core.windows.net'):
            logger.error(f"URL does not appear to be an Azure Blob Storage URL: {storage_url}")
            return None

        # Extract path components - remove leading slash
        path = parsed_url.path
        if path.startswith('/'):
            path = path[1:]

        # Split into container and blob name
        parts = path.split('/', 1)
        if len(parts) != 2:
            logger.error(f"URL does not contain both container and blob name: {storage_url}")
            return None

        container_name = parts[0]
        blob_name = parts[1]

        logger.info(f"Parsed URL - Container: {container_name}, Blob: {blob_name}")

        # Create a BlobServiceClient using the connection string
        blob_service_client = BlobServiceClient.from_connection_string(connection_string)

        # Get the container client
        container_client = blob_service_client.get_container_client(container_name)

        # Get the blob client
        blob_client = container_client.get_blob_client(blob_name)

        # Create a temporary file to store the downloaded PowerPoint
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
            temp_path = temp_file.name

            # Download the blob content
            blob_data = blob_client.download_blob()

            # Write the PowerPoint content to the temporary file
            temp_file.write(blob_data.readall())

        # Process the PowerPoint file
        documents = extract_text_from_powerpoint(temp_path)

        # Clean up the temporary file
        os.unlink(temp_path)

        # Add source metadata to documents
        if documents:
            for doc in documents:
                if not doc.metadata:
                    doc.metadata = {}
                doc.metadata["source"] = storage_url

        logger.info(f"Successfully loaded PowerPoint from storage URL: {storage_url}")
        return documents

    except Exception as e:
        logger.error(f"Error loading PowerPoint from storage URL {storage_url}: {str(e)}")
        # Clean up temp file if it exists
        if 'temp_path' in locals():
            try:
                os.unlink(temp_path)
            except:
                pass
        return None


def load_pptx_from_url(pptx_url: str) -> Optional[List[Document]]:
    """
    Downloads a PowerPoint file from a URL and extracts its content.

    Args:
        pptx_url: URL of the PowerPoint file to download and process

    Returns:
        List[Document]: List of Document objects or None if loading fails
    """
    try:
        # Set up headers to mimic a regular browser
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,image/apng,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.9',
            'Accept-Encoding': 'gzip, deflate, br',
            'Connection': 'keep-alive',
            'Upgrade-Insecure-Requests': '1',
            'Cache-Control': 'max-age=0'
        }

        # Download the PowerPoint with headers
        response = requests.get(pptx_url, stream=True, headers=headers)
        response.raise_for_status()  # Raise exception for HTTP errors

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
            temp_path = temp_file.name
            # Write the PowerPoint content to the temporary file
            for chunk in response.iter_content(chunk_size=8192):
                temp_file.write(chunk)

        # Process the PowerPoint file
        documents = extract_text_from_powerpoint(temp_path)

        # Clean up the temporary file
        os.unlink(temp_path)

        # Add source metadata to documents
        if documents:
            for doc in documents:
                if not doc.metadata:
                    doc.metadata = {}
                doc.metadata["source"] = pptx_url

        logger.info(f"Successfully loaded PowerPoint from {pptx_url}")
        return documents

    except requests.RequestException as e:
        logger.error(f"Failed to download PowerPoint from {pptx_url}: {str(e)}")
        return None
    except Exception as e:
        logger.error(f"Error loading PowerPoint from {pptx_url}: {str(e)}")
        # Clean up temp file if it exists
        if 'temp_path' in locals():
            try:
                os.unlink(temp_path)
            except:
                pass
        return None


def load_pptx_from_private_url(
        private_url: str,
        connection_string: str = None,
        container_name: str = None
) -> Optional[List[Document]]:
    """
    Parses a private Azure URL and loads the PowerPoint from Azure Blob Storage.

    Args:
        private_url: Private URL pointing to Azure Blob (could be in various formats)
        connection_string: Azure Storage connection string
        container_name: Name of the Azure Blob container (can be overridden if contained in URL)

    Returns:
        List[Document]: List of Document objects or None if loading fails
    """
    try:
        # For full storage URLs, use the direct storage URL loader
        if private_url.startswith('https://') and '.blob.core.windows.net' in private_url:
            return load_pptx_from_storage_url(
                storage_url=private_url,
                connection_string=connection_string
            )

        # If it's a regular HTTP URL, use the standard URL loader
        if private_url.startswith('http://') or private_url.startswith('https://'):
            if '.blob.core.windows.net' not in private_url:
                return load_pptx_from_url(private_url)

        # Use environment variables if not provided
        if connection_string is None:
            connection_string = os.getenv("AZURE_STORAGE_CONNECTION_STRING")
            if not connection_string:
                raise ValueError("Azure connection string not provided and not found in environment variables")

        if container_name is None:
            container_name = os.getenv("AZURE_STORAGE_CONTAINER_NAME")
            if not container_name:
                raise ValueError("Container name not provided and not found in environment variables")

        # Extract the blob name from the private URL
        if private_url.startswith("https://"):
            # Parse URL to extract blob path
            parts = private_url.split('/')
            if len(parts) < 5:
                logger.error(f"Invalid Azure Blob URL format: {private_url}")
                return None

            # If URL contains container name, use it instead of the provided one
            url_container = parts[3]
            if container_name and url_container != container_name:
                logger.info(f"Using container name from URL: {url_container} instead of provided: {container_name}")
                container_name = url_container

            # Get the blob path (everything after the container name)
            blob_name = '/'.join(parts[4:])
            if not blob_name:
                logger.error(f"Could not extract blob name from URL: {private_url}")
                return None
        else:
            # Assume the private_url is just the blob name
            blob_name = private_url

        logger.info(f"Extracted blob name: {blob_name} from URL: {private_url}")

        # Create a BlobServiceClient using the connection string
        blob_service_client = BlobServiceClient.from_connection_string(connection_string)

        # Get the container client
        container_client = blob_service_client.get_container_client(container_name)

        # Get the blob client
        blob_client = container_client.get_blob_client(blob_name)

        # Create a temporary file to store the downloaded PowerPoint
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
            temp_path = temp_file.name

            # Download the blob content
            blob_data = blob_client.download_blob()

            # Write the PowerPoint content to the temporary file
            temp_file.write(blob_data.readall())

        # Process the PowerPoint file
        documents = extract_text_from_powerpoint(temp_path)

        # Clean up the temporary file
        os.unlink(temp_path)

        # Add source metadata to documents
        if documents:
            for doc in documents:
                if not doc.metadata:
                    doc.metadata = {}
                doc.metadata["source"] = private_url

        logger.info(f"Successfully loaded PowerPoint from private URL: {private_url}")
        return documents

    except Exception as e:
        logger.error(f"Error loading PowerPoint from private URL {private_url}: {str(e)}")
        # Clean up temp file if it exists
        if 'temp_path' in locals():
            try:
                os.unlink(temp_path)
            except:
                pass
        return None


def extract_text_from_powerpoint(file_path: str) -> Optional[List[Document]]:
    """
    Extracts text and content from a PowerPoint file.

    Args:
        file_path: Path to the PowerPoint file

    Returns:
        List[Document]: List of Document objects or None if extraction fails
    """
    try:
        # Load the presentation
        presentation = Presentation(file_path)
        documents = []

        # Extract text from each slide
        for slide_index, slide in enumerate(presentation.slides):
            slide_text = ""
            slide_notes = ""

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += shape.text + "\n"

            # Extract notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for note_shape in notes_slide.shapes:
                    if hasattr(note_shape, "text") and note_shape.text:
                        slide_notes += note_shape.text + "\n"

            # Create a Document object for this slide
            if slide_text.strip() or slide_notes.strip():
                content = f"Slide {slide_index + 1}:\n{slide_text.strip()}"
                if slide_notes.strip():
                    content += f"\n\nSlide Notes:\n{slide_notes.strip()}"

                doc = Document(
                    page_content=content,
                    metadata={
                        "source": file_path,
                        "slide_number": slide_index + 1,
                        "type": "powerpoint_slide"
                    }
                )
                documents.append(doc)

        logger.info(f"Successfully extracted {len(documents)} slides from PowerPoint file: {file_path}")
        return documents

    except Exception as e:
        logger.error(f"Error extracting text from PowerPoint file {file_path}: {str(e)}")
        return None


def is_powerpoint_url(url: str) -> bool:
    """
    Determines if a URL is likely pointing to a PowerPoint file

    Args:
        url: URL to check

    Returns:
        bool: True if the URL is likely a PowerPoint file, False otherwise
    """
    return url.lower().endswith(
        ('.ppt', '.pptx', '.pps', '.ppsx')) or '/ppt' in url.lower() or 'presentation' in url.lower()




# ------- This is AI powerpoint presentation -----------



class SlideExtract(BaseModel):
    """Model for extracted slide content"""
    slide_number: int
    title: str
    content: str
    notes: Optional[str] = None


class PresentationExtract(BaseModel):
    """Model for extracted presentation content"""
    title: str
    slides: List[SlideExtract]
    total_slides: int


def process_powerpoint_file(file_content: bytes) -> PresentationExtract:
    """
    Process PowerPoint file bytes and extract slide content

    Args:
        file_content: Raw bytes of the PowerPoint file

    Returns:
        PresentationExtract: Structured content from the presentation
    """
    try:
        # Create a temporary file to store the PowerPoint
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
            temp_path = temp_file.name
            temp_file.write(file_content)

        # Extract content from the PowerPoint file
        presentation = Presentation(temp_path)
        slides = []
        presentation_title = "Untitled Presentation"

        # Try to extract presentation title from properties
        if hasattr(presentation.core_properties, 'title') and presentation.core_properties.title:
            presentation_title = presentation.core_properties.title

        # Extract text from each slide
        for slide_idx, slide in enumerate(presentation.slides):
            slide_title = ""
            slide_content = ""
            slide_notes = ""

            # Extract title from first shape if it appears to be a title
            if slide.shapes.title:
                slide_title = slide.shapes.title.text

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # Skip if this is the title we already extracted
                    if shape.text == slide_title:
                        continue
                    slide_content += shape.text + "\n"

            # Extract notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for note_shape in notes_slide.shapes:
                    if hasattr(note_shape, "text") and note_shape.text:
                        slide_notes += note_shape.text + "\n"

            # Create a SlideExtract object
            slide_extract = SlideExtract(
                slide_number=slide_idx + 1,
                title=slide_title,
                content=slide_content.strip(),
                notes=slide_notes.strip() if slide_notes.strip() else None
            )
            slides.append(slide_extract)

        # Clean up the temporary file
        os.unlink(temp_path)

        return PresentationExtract(
            title=presentation_title,
            slides=slides,
            total_slides=len(slides)
        )

    except Exception as e:
        logger.error(f"Error processing PowerPoint file: {str(e)}")
        # Clean up temp file if it exists
        if 'temp_path' in locals():
            try:
                os.unlink(temp_path)
            except:
                pass
        raise

async def process_powerpoint_from_url(powerpoint_url: str) -> PresentationExtract:
    """
    Process PowerPoint from Azure Blob Storage URL

    Args:
        powerpoint_url: Azure Blob Storage URL to the PowerPoint file

    Returns:
        PresentationExtract: Structured content from the presentation
    """
    try:
        # Check if the URL is an Azure Blob Storage URL
        if 'blob.core.windows.net' in powerpoint_url:
            # Get the PowerPoint file content from Azure Blob Storage
            file_content = await download_from_azure_blob(powerpoint_url)
        else:
            # If it's a regular HTTP URL, use the standard HTTP download
            file_content = await download_from_http(powerpoint_url)

        if not file_content:
            raise ValueError(f"Failed to download PowerPoint from URL: {powerpoint_url}")

        # Process the PowerPoint file
        return await process_powerpoint_content(file_content)
    except Exception as e:
        logger.error(f"Error processing PowerPoint from URL {powerpoint_url}: {str(e)}")
        raise


async def download_from_azure_blob(blob_url: str) -> bytes:
    """
    Download file from Azure Blob Storage using the URL

    Args:
        blob_url: Azure Blob Storage URL

    Returns:
        bytes: File content
    """
    try:
        # Parse the Azure Blob Storage URL
        parsed_url = urlparse(blob_url)

        # Extract account name from the host
        # Format: accountname.blob.core.windows.net
        account_name = parsed_url.netloc.split('.')[0]

        # Extract container name and blob path from the URL path
        path_parts = parsed_url.path.strip('/').split('/', 1)
        if len(path_parts) != 2:
            raise ValueError(f"Invalid Azure Blob URL format: {blob_url}")

        container_name = path_parts[0]
        blob_name = path_parts[1]

        logger.info(
            f"Downloading from Azure Blob Storage: Account={account_name}, Container={container_name}, Blob={blob_name}")

        # Get connection string from environment or use SAS token from URL
        connection_string = os.getenv("AZURE_STORAGE_CONNECTION_STRING")

        if connection_string:
            # Use connection string if available
            blob_service_client = BlobServiceClient.from_connection_string(connection_string)
            container_client = blob_service_client.get_container_client(container_name)
            blob_client = container_client.get_blob_client(blob_name)

            # Download the blob content
            download_stream = await asyncio.to_thread(blob_client.download_blob)
            file_content = await asyncio.to_thread(download_stream.readall)
        else:
            # If no connection string, assume the URL contains a SAS token
            # Use anonymous access with the full URL
            blob_service_client = BlobServiceClient(account_url=f"https://{account_name}.blob.core.windows.net")
            container_client = blob_service_client.get_container_client(container_name)
            blob_client = container_client.get_blob_client(blob_name)

            # The URL should contain the SAS token for authentication
            download_stream = await asyncio.to_thread(
                lambda: blob_client.download_blob(credential=parsed_url.query)
            )
            file_content = await asyncio.to_thread(download_stream.readall)

        logger.info(f"Successfully downloaded PowerPoint from Azure Blob Storage, size: {len(file_content)} bytes")
        return file_content

    except Exception as e:
        logger.error(f"Error downloading from Azure Blob Storage: {str(e)}")
        raise


async def process_powerpoint_content(file_content: bytes) -> PresentationExtract:
    """
    Process PowerPoint file content and extract slide information

    Args:
        file_content: Raw bytes of the PowerPoint file

    Returns:
        PresentationExtract: Structured content from the presentation
    """
    try:
        # Create a temporary file to store the PowerPoint
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
            temp_path = temp_file.name
            temp_file.write(file_content)

        # Process the PowerPoint file
        try:
            # Load the presentation
            presentation = await asyncio.to_thread(Presentation, temp_path)

            slides = []
            presentation_title = "Untitled Presentation"

            # Try to extract presentation title from properties
            if hasattr(presentation.core_properties, 'title') and presentation.core_properties.title:
                presentation_title = presentation.core_properties.title

            # Extract text from each slide
            for slide_idx, slide in enumerate(presentation.slides):
                slide_title = ""
                slide_content = ""
                slide_notes = ""

                # Extract title from first shape if it appears to be a title
                if slide.shapes.title:
                    slide_title = slide.shapes.title.text

                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Skip if this is the title we already extracted
                        if shape.text == slide_title:
                            continue
                        slide_content += shape.text + "\n"

                # Extract notes
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    for note_shape in notes_slide.shapes:
                        if hasattr(note_shape, "text") and note_shape.text:
                            slide_notes += note_shape.text + "\n"

                # Create a SlideExtract object
                slide_extract = SlideExtract(
                    slide_number=slide_idx + 1,
                    title=slide_title,
                    content=slide_content.strip(),
                    notes=slide_notes.strip() if slide_notes.strip() else None
                )
                slides.append(slide_extract)

            # Create the presentation extract
            presentation_extract = PresentationExtract(
                title=presentation_title,
                slides=slides,
                total_slides=len(slides)
            )

            return presentation_extract

        finally:
            # Clean up the temporary file
            os.unlink(temp_path)

    except Exception as e:
        logger.error(f"Error processing PowerPoint content: {str(e)}")
        raise


def convert_to_setup_format(
        presentation_extract: PresentationExtract,
        session_id: str,
        presenter_ranges: List[Tuple[int, int]],
        student_persona: Optional[Dict] = None
) -> Dict:
    """
    Convert extracted presentation to setup format with both list and dict formats

    Args:
        presentation_extract: Extracted presentation content
        session_id: Session identifier
        presenter_ranges: List of tuples with ranges (start, end) for human presenter
        student_persona: Optional customizations for the AI persona

    Returns:
        Dict: Formatted presentation setup data
    """
    # Format slide contents for the API in both list and dict formats
    slide_contents_list = []
    slide_contents_dict = {}

    for slide in presentation_extract.slides:
        content = f"{slide.title}\n\n{slide.content}"
        if slide.notes:
            content += f"\n\nNotes: {slide.notes}"

        # Add to list format
        slide_contents_list.append({
            "slide_number": slide.slide_number,
            "content": content
        })

        # Add to dict format (key is the slide number)
        slide_contents_dict[slide.slide_number] = content

    # Format presenter ranges
    presenter_range_objects = []
    for start, end in presenter_ranges:
        presenter_range_objects.append({
            "start": start,
            "end": end
        })

    # Create setup data
    setup_data = {
        "session_id": session_id,
        "slide_contents": slide_contents_list,
        "slide_contents_dict": slide_contents_dict,
        "presenter_ranges": presenter_range_objects,
        "presentation_title": presentation_extract.title,
        "presentation_topic": presentation_extract.title  # Use title as topic if not provided
    }

    # Add student persona if provided
    if student_persona:
        setup_data["student_persona"] = student_persona

    return setup_data

async def download_from_http(url: str) -> bytes:
    """
    Download file from a standard HTTP URL

    Args:
        url: HTTP URL

    Returns:
        bytes: File content
    """
    try:
        import aiohttp

        async with aiohttp.ClientSession() as session:
            async with session.get(url) as response:
                if response.status != 200:
                    raise ValueError(f"Failed to download file, status code: {response.status}")

                file_content = await response.read()

        logger.info(f"Successfully downloaded PowerPoint from HTTP URL, size: {len(file_content)} bytes")
        return file_content

    except Exception as e:
        logger.error(f"Error downloading from HTTP URL: {str(e)}")
        raise

async def process_powerpoint_from_upload(file_data: bytes) -> PresentationExtract:
    """
    Process PowerPoint from uploaded file data

    Args:
        file_data: Raw bytes from uploaded file

    Returns:
        PresentationExtract: Structured content from the presentation
    """
    return process_powerpoint_file(file_data)


async def process_powerpoint_from_base64(base64_data: str) -> PresentationExtract:
    """
    Process PowerPoint from base64-encoded string

    Args:
        base64_data: Base64-encoded PowerPoint data

    Returns:
        PresentationExtract: Structured content from the presentation
    """
    # Remove data URI prefix if present
    if base64_data.startswith('data:'):
        # Extract just the base64 part
        base64_data = base64_data.split(',')[1]

    # Decode base64 to bytes
    file_content = base64.b64decode(base64_data)

    return process_powerpoint_file(file_content)


def convert_to_presentation_setup_enhanced(
        presentation_extract: PresentationExtract,
        session_id: str,
        presenter_ranges: List[Tuple[int, int]],
        student_persona: Optional[Dict] = None
) -> Dict:
    """
    Convert extracted presentation to setup format with enhanced validation

    Args:
        presentation_extract: Extracted presentation content
        session_id: Session identifier
        presenter_ranges: List of tuples with ranges (start, end) for human presenter
        student_persona: Optional customizations for the AI persona

    Returns:
        Dict: Formatted presentation setup data with both list and dict formats
    """
    # Format slide contents for the API
    slide_contents_list = []
    slide_contents_dict = {}

    for slide in presentation_extract.slides:
        content = f"{slide.title}\n\n{slide.content}"
        if slide.notes:
            content += f"\n\nNotes: {slide.notes}"

        # Add to list format
        slide_contents_list.append({
            "slide_number": slide.slide_number,
            "content": content
        })

        # Add to dict format
        slide_contents_dict[slide.slide_number] = content

    # Format presenter ranges
    presenter_range_objects = []
    for start, end in presenter_ranges:
        presenter_range_objects.append({
            "start": start,
            "end": end
        })

    # Create setup data
    setup_data = {
        "session_id": session_id,
        "slide_contents": slide_contents_list,
        "slide_contents_dict": slide_contents_dict,  # Also include dict format
        "presenter_ranges": presenter_range_objects,
        "presentation_title": presentation_extract.title,
        "presentation_topic": presentation_extract.title  # Use title as topic if not provided
    }

    # Add student persona if provided
    if student_persona:
        setup_data["student_persona"] = student_persona

    return setup_data


# Add this function to help debug slide content structures
def validate_slide_content(slide_content) -> bool:
    """
    Validate slide content structure to ensure it's in the correct format

    Args:
        slide_content: Slide content to validate

    Returns:
        bool: True if valid, False otherwise
    """
    try:
        if isinstance(slide_content, list):
            # Check list format - each item must have slide_number and content
            for item in slide_content:
                if not isinstance(item, dict):
                    return False
                if "slide_number" not in item or "content" not in item:
                    return False
            return True
        elif isinstance(slide_content, dict):
            # Check dict format - keys must be slide numbers
            for key in slide_content.keys():
                if not isinstance(key, int):
                    return False
            return True
        else:
            return False
    except Exception:
        return False